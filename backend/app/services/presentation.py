"""论文分享 PPT 生成 v2：显式几何排版 harness（借鉴 Claude Code 制 PPT 工作流）。

原则：
- 模板只用品牌页（封面/目录/节标题/结尾）与「仅标题」的标题条；正文一律显式几何排版
  （卡片网格/彩头分栏/对比/表格/图文），杜绝占位符裁切与大片留白；
- 设计 token：固定色板 + 字号四档（30/20/18/16）+ 中文字体（含 a:ea 东亚字体强制）；
- 图片 PIL 读宽高比等比适配；宽图上文下图、窄图左文右图；
- 「lead｜正文」约定：要点内 ｜ 前的引导词渲染为 20pt 蓝色加粗；
- validate_deck_spec 确定性校验（短标题/无破折号/密度上下限）供反馈迭代。
"""

import re
import shutil
import subprocess
import tempfile
from io import BytesIO
from pathlib import Path
from typing import Literal

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pydantic import BaseModel, ConfigDict, Field

TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "assets" / "presentation_template.pptx"

T1, T2, T3, T4 = 30, 20, 18, 16  # 大标题/小标题/正文/小字
FONT = "微软雅黑"

BLUE = RGBColor(0x15, 0x3F, 0x82)
ORANGE = RGBColor(0xEE, 0x82, 0x2F)
TEAL = RGBColor(0x30, 0xC0, 0xB4)
GREEN = RGBColor(0x75, 0xBD, 0x42)
RED = RGBColor(0xE5, 0x4C, 0x5E)
GREY = RGBColor(0x59, 0x59, 0x59)
DARK = RGBColor(0x33, 0x33, 0x33)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COL_COLORS = [BLUE, TEAL, ORANGE, GREEN]

PAGE_W, PAGE_H = 13.33, 7.5
MX = 0.55  # 左右边距
BODY_Y = 1.15  # 标题条下正文起点

TITLE_MAX_CHARS = 20
_DASH_RE = re.compile(r"—|–|——|--")

_LAYOUT = {"cover": 0, "toc": 1, "section": 2, "blank": 4, "closing": 9}


class DeckItem(BaseModel):
    model_config = ConfigDict(extra="forbid")

    title: str = ""
    body: str = ""


class DeckColumn(BaseModel):
    model_config = ConfigDict(extra="forbid")

    header: str = ""
    bullets: list[str] = Field(default_factory=list)


class DeckTable(BaseModel):
    model_config = ConfigDict(extra="forbid")

    headers: list[str] = Field(default_factory=list, max_length=6)
    rows: list[list[str]] = Field(default_factory=list, max_length=8)


class DeckSlide(BaseModel):
    model_config = ConfigDict(extra="forbid")

    kind: Literal[
        "cover",
        "toc",
        "section",
        "content",
        "cards",
        "columns",
        "compare",
        "table",
        "figure",
        "closing",
    ]
    title: str = ""
    subtitle: str = ""
    presenter: str = ""
    date: str = ""
    items: list[str] = Field(default_factory=list)  # toc
    bullets: list[str] = Field(default_factory=list)  # content；"lead｜正文"；行首 "- " 为二级
    cards: list[DeckItem] = Field(default_factory=list, max_length=6)  # cards：两列卡片网格
    columns: list[DeckColumn] = Field(default_factory=list, max_length=4)  # columns：彩头分栏
    left: DeckColumn | None = None  # compare
    right: DeckColumn | None = None
    table: DeckTable | None = None
    figure_index: int | None = None
    caption: str = ""  # figure 讲解（16pt）
    summary: str = ""  # 页底总结条（"lead｜正文"）
    notes: str = ""


class DeckSpec(BaseModel):
    model_config = ConfigDict(extra="forbid")

    title: str
    slides: list[DeckSlide] = Field(min_length=3, max_length=25)


# ---- 确定性校验 ----


def _texts_of(s: DeckSlide) -> list[str]:
    out = [s.title, s.subtitle, s.caption, s.summary, *s.bullets, *s.items]
    out += [c.title + c.body for c in s.cards]
    for col in [*s.columns, s.left, s.right]:
        if col:
            out += [col.header, *col.bullets]
    if s.table:
        out += s.table.headers + [c for row in s.table.rows for c in row]
    return [t for t in out if t]


def validate_deck_spec(spec: DeckSpec, *, figure_indices: set[int]) -> list[str]:
    errors: list[str] = []
    for i, s in enumerate(spec.slides):
        where = f"第 {i + 1} 页「{s.title or s.kind}」"
        for t in _texts_of(s):
            if _DASH_RE.search(t):
                errors.append(f"{where}文本「{t[:16]}…」含破折号（规范禁止，用逗号/冒号/、改写）")
        if s.title and len(s.title) > TITLE_MAX_CHARS:
            errors.append(f"{where}标题 {len(s.title)} 字超上限 {TITLE_MAX_CHARS}（要短标题）")
        if s.kind == "content":
            if not (2 <= len(s.bullets) <= 6):
                errors.append(
                    f"{where}要点 {len(s.bullets)} 条（content 页需 2-6 条，太少显空、太多拆页）"
                )
            for b in s.bullets:
                if len(b.lstrip("- ")) > 60:
                    errors.append(f"{where}要点「{b[:16]}…」超 60 字（引导词｜短句，细节放 notes）")
        if s.kind == "cards" and not (2 <= len(s.cards) <= 6):
            errors.append(f"{where}卡片 {len(s.cards)} 张（需 2-6 张）")
        for c in s.cards:
            if len(c.body) > 90:
                errors.append(f"{where}卡片「{c.title}」正文超 90 字")
        if s.kind == "columns" and not (2 <= len(s.columns) <= 4):
            errors.append(f"{where}分栏 {len(s.columns)} 栏（需 2-4 栏）")
        for col in [*s.columns, s.left, s.right]:
            if col:
                for b in col.bullets:
                    if len(b) > 44:
                        errors.append(f"{where}栏目要点「{b[:14]}…」超 44 字")
        if s.kind == "table" and (not s.table or not s.table.headers or not s.table.rows):
            errors.append(f"{where}table 页缺 headers/rows")
        if s.table:
            for row in s.table.rows:
                for c in row:
                    if len(c) > 30:
                        errors.append(f"{where}表格单元「{c[:12]}…」超 30 字")
        if s.kind == "toc" and not (3 <= len(s.items) <= 5):
            errors.append(f"{where}目录需 3-5 条")
        if s.kind == "figure":
            if s.figure_index is None or s.figure_index not in figure_indices:
                errors.append(f"{where}figure_index 非法（可用：{sorted(figure_indices)}）")
            if len(s.caption) < 15:
                errors.append(f"{where}图片讲解太短（要说清画了什么、支撑什么论点、看图中哪里）")
    if spec.slides[0].kind != "cover":
        errors.append("第 1 页必须是 cover")
    if figure_indices and not any(s.kind == "figure" for s in spec.slides):
        errors.append("有可用配图但没有 figure 页（要求图文并茂，可用配图尽量用上）")
    return errors


def _sv(text: str) -> str:
    return _DASH_RE.sub("，", text)


# ---- 几何排版原语（borrowed from Claude Code PPT harness） ----


def _set_font(run, size: int, bold: bool = False, color: RGBColor | None = DARK) -> None:
    f = run.font
    f.size, f.bold, f.name = Pt(size), bold, FONT
    if color is not None:  # None = 继承模板配色（品牌页标题用，防蓝字叠蓝底隐形）
        f.color.rgb = color
    rPr = run._r.get_or_add_rPr()  # noqa: SLF001
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def _para(tf, first: bool):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def _write(tf, paras, anchor=MSO_ANCHOR.TOP, ml=0.1, mt=0.06) -> None:
    """paras: [{runs: [(text,size,bold,color)], align?, space_after?, line?}]"""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(ml)
    tf.margin_top = tf.margin_bottom = Inches(mt)
    for i, p in enumerate(paras):
        para = _para(tf, i == 0)
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("space_after") is not None:
            para.space_after = Pt(p["space_after"])
        if p.get("line") is not None:
            para.line_spacing = p["line"]
        for t, size, bold, color in p["runs"]:
            run = para.add_run()
            run.text = _sv(t)
            _set_font(run, size, bold, color)


def _text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _write(tb.text_frame, paras, anchor=anchor, ml=0, mt=0)
    return tb


def _box(slide, x, y, w, h, fill=LIGHT, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.adjustments[0] = 0.08
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.2)
    shp.shadow.inherit = False
    return shp


def _fit_picture(slide, data: bytes, box_x, box_y, box_w, box_h):
    with Image.open(BytesIO(data)) as im:
        iw, ih = im.size
    ar = iw / max(ih, 1)
    w = box_w
    h = w / ar
    if h > box_h:
        h = box_h
        w = h * ar
    x = box_x + (box_w - w) / 2
    y = box_y + (box_h - h) / 2
    return slide.shapes.add_picture(BytesIO(data), Inches(x), Inches(y), Inches(w), Inches(h))


def _lead_runs(text: str, *, lead_color=BLUE, size=T3):
    """「lead｜正文」→ lead 20pt 加粗彩色 + 正文；无 ｜ 则整句正文。"""
    if "｜" in text[:14]:
        lead, rest = text.split("｜", 1)
        return [(lead + "｜", T2, True, lead_color), (rest, size, False, DARK)]
    return [(text, size, False, DARK)]


def _summary_bar(slide, text: str, y=6.45) -> None:
    bar = _box(slide, MX, y, PAGE_W - 2 * MX, 0.78, fill=None, line=BLUE)
    _write(
        bar.text_frame,
        [{"runs": _lead_runs(text), "align": PP_ALIGN.CENTER}],
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _ph(slide, idx: int):
    for p in slide.placeholders:
        if p.placeholder_format.idx == idx:
            return p
    return None


def _set_ph(slide, idx: int, text: str, size: int, bold: bool = False) -> None:
    ph = _ph(slide, idx)
    if ph is None:
        return
    ph.text = _sv(text)
    for para in ph.text_frame.paragraphs:
        for run in para.runs:
            _set_font(
                run,
                size,
                bold,
                color=run.font.color.rgb if run.font.color and run.font.color.type else DARK,
            )


def _drop_template_slides(prs: Presentation) -> None:
    sld_lst = prs.slides._sldIdLst  # noqa: SLF001
    for sld in list(sld_lst):
        prs.part.drop_rel(sld.rId)
        sld_lst.remove(sld)


# ---- 各版式渲染 ----


def _render_content(slide, s: DeckSlide) -> None:
    n = len(s.bullets)
    paras = []
    for b in s.bullets:
        sub = b.startswith("- ")
        t = b[2:] if sub else b
        runs = [("▪ ", T4, False, GREY), (t, T4, False, GREY)] if sub else _lead_runs(t)
        paras.append({"runs": runs, "space_after": 14 if n <= 4 else 9, "line": 1.15})
    h = 5.1 if not s.summary else 5.0
    _text(
        slide, MX + 0.15, BODY_Y + 0.15, PAGE_W - 2 * MX - 0.3, h, paras, anchor=MSO_ANCHOR.MIDDLE
    )


def _render_cards(slide, s: DeckSlide) -> None:
    n = len(s.cards)
    rows = (n + 1) // 2
    area_h = (6.3 if not s.summary else 5.15) - BODY_Y
    ch = min(2.1, (area_h - 0.15 * (rows - 1)) / rows)
    cw = (PAGE_W - 2 * MX - 0.35) / 2
    for i, c in enumerate(s.cards):
        x = MX + (i % 2) * (cw + 0.35)
        y = BODY_Y + 0.1 + (i // 2) * (ch + 0.15)
        card = _box(slide, x, y, cw, ch, fill=LIGHT)
        _write(
            card.text_frame,
            [
                {"runs": [(c.title, T2, True, BLUE)], "space_after": 4},
                {"runs": [(c.body, T3 if len(c.body) <= 60 else T4, False, DARK)], "line": 1.12},
            ],
            anchor=MSO_ANCHOR.MIDDLE,
            ml=0.16,
        )


def _render_columns(slide, s: DeckSlide, cols: list[DeckColumn], *, vs: bool = False) -> None:
    n = len(cols)
    gap = 0.22
    cw = (PAGE_W - 2 * MX - gap * (n - 1)) / n
    body_h = (6.3 if not s.summary else 5.15) - BODY_Y - 0.65
    for i, col in enumerate(cols):
        x = MX + i * (cw + gap)
        color = RED if vs and i == 0 else TEAL if vs else COL_COLORS[i % len(COL_COLORS)]
        head = _box(slide, x, BODY_Y + 0.1, cw, 0.5, fill=color)
        _write(
            head.text_frame,
            [{"runs": [(col.header, T2, True, WHITE)], "align": PP_ALIGN.CENTER}],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        body = _box(slide, x, BODY_Y + 0.68, cw, body_h, fill=LIGHT)
        paras = [
            {"runs": [("• " + b, T4, False, DARK)], "line": 1.12, "space_after": 8}
            for b in col.bullets
        ]
        _write(body.text_frame, paras, ml=0.14, mt=0.12)


def _render_table(slide, s: DeckSlide) -> None:
    tbl_def = s.table
    assert tbl_def is not None
    ncols = len(tbl_def.headers)
    nrows = len(tbl_def.rows) + 1
    height = min(5.0, 0.5 + 0.45 * (nrows - 1))
    shape = slide.shapes.add_table(
        nrows, ncols, Inches(MX), Inches(BODY_Y + 0.15), Inches(PAGE_W - 2 * MX), Inches(height)
    )
    tbl = shape.table
    for ri in range(nrows):
        cells = tbl_def.headers if ri == 0 else tbl_def.rows[ri - 1]
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ri == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = _sv(cells[ci] if ci < len(cells) else "")
            _set_font(run, T4, ri == 0, WHITE if ri == 0 else DARK)


def _render_figure(slide, s: DeckSlide, figures: dict[int, bytes]) -> None:
    data = figures.get(s.figure_index) if s.figure_index is not None else None
    cap_paras = [{"runs": [(s.caption, T4, False, DARK)], "line": 1.25}]
    if data is None:
        _text(slide, MX, BODY_Y + 0.2, PAGE_W - 2 * MX, 1.0, cap_paras)
        return
    with Image.open(BytesIO(data)) as im:
        ar = im.size[0] / max(im.size[1], 1)
    if ar >= 1.7:  # 宽图：上讲解、下大图
        _text(slide, MX + 0.1, BODY_Y + 0.05, PAGE_W - 2 * MX - 0.2, 0.9, cap_paras)
        _fit_picture(slide, data, MX, BODY_Y + 1.0, PAGE_W - 2 * MX, 7.1 - (BODY_Y + 1.0))
    else:  # 窄图：左讲解、右图
        _text(slide, MX + 0.1, BODY_Y + 0.2, 5.5, 5.4, cap_paras, anchor=MSO_ANCHOR.MIDDLE)
        _fit_picture(slide, data, 6.35, BODY_Y + 0.1, PAGE_W - MX - 6.35, 5.9)


def build_deck(spec: DeckSpec, figures: dict[int, bytes]) -> bytes:
    prs = Presentation(str(TEMPLATE_PATH))
    _drop_template_slides(prs)

    for s in spec.slides:
        template_kind = s.kind if s.kind in _LAYOUT else "blank"
        slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT[template_kind]])
        if slide.shapes.title is not None and s.title:
            slide.shapes.title.text = _sv(s.title)
            # 品牌页（封面/节标题等）标题继承模板配色与位置；只统一字号与字体
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    _set_font(run, T1, True, color=None)

        if s.kind == "cover":
            _set_ph(slide, 1, s.subtitle, T2)
            _set_ph(slide, 20, s.presenter or "汇报人：（待填）", T4)
            _set_ph(slide, 21, s.date, T4)
        elif s.kind == "toc":
            for i in range(5):
                num_ph, item_ph = _ph(slide, 21 + i), _ph(slide, 11 + i)
                if i < len(s.items):
                    _set_ph(slide, 21 + i, f"{i + 1:02d}", T2, bold=True)
                    _set_ph(slide, 11 + i, s.items[i], T3)
                else:  # 多余目录占位符整对删除，避免空框
                    for ph in (num_ph, item_ph):
                        if ph is not None:
                            ph._element.getparent().remove(ph._element)  # noqa: SLF001
        elif s.kind == "section":
            _set_ph(slide, 1, s.subtitle or s.title, T2)
        elif s.kind == "closing":
            _set_ph(slide, 13, s.subtitle, T2)
        elif s.kind == "content":
            _render_content(slide, s)
        elif s.kind == "cards":
            _render_cards(slide, s)
        elif s.kind == "columns":
            _render_columns(slide, s, s.columns)
        elif s.kind == "compare":
            cols = [c for c in (s.left, s.right) if c]
            _render_columns(slide, s, cols, vs=True)
        elif s.kind == "table":
            _render_table(slide, s)
        elif s.kind == "figure":
            _render_figure(slide, s, figures)

        if s.summary and s.kind in ("content", "cards", "columns", "compare", "table"):
            _summary_bar(slide, s.summary)
        if s.notes:
            slide.notes_slide.notes_text_frame.text = _sv(s.notes)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---- 渲染反馈（soffice → pdf → PNG） ----


def soffice_available() -> bool:
    return shutil.which("soffice") is not None


def render_slide_images(pptx_bytes: bytes, *, max_pages: int = 12) -> list[bytes]:
    if not soffice_available():
        return []
    import fitz

    with tempfile.TemporaryDirectory(prefix="polaris-pptx-") as tmp:
        src = Path(tmp) / "deck.pptx"
        src.write_bytes(pptx_bytes)
        try:
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, str(src)],
                capture_output=True,
                timeout=120,
                check=True,
            )
        except (subprocess.SubprocessError, OSError):
            return []
        pdf = Path(tmp) / "deck.pdf"
        if not pdf.exists():
            return []
        images: list[bytes] = []
        with fitz.open(pdf) as doc:
            for page in doc.pages(0, min(max_pages, doc.page_count)):
                images.append(page.get_pixmap(dpi=110).tobytes("png"))
        return images
