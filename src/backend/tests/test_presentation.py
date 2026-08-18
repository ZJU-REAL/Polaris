"""论文分享 PPT：甲板校验/模板渲染/SKILL.md 导入/任务创建。"""

import io
import uuid

import pytest
from pptx import Presentation as PptxFile
from pydantic import ValidationError

from app.services.presentation import DeckSpec, build_deck, validate_deck_spec
from tests.conftest import add_paper, register_and_login

DECK = {
    "title": "测试分享",
    "slides": [
        {
            "kind": "cover",
            "title": "自我奖励语言模型",
            "subtitle": "arXiv 2024",
            "presenter": "汇报人：测试",
        },
        {"kind": "toc", "title": "目录", "items": ["背景", "方法", "结果"]},
        {
            "kind": "content",
            "title": "研究背景",
            "bullets": ["奖励模型受限于人类偏好数据", "- 质量随规模饱和"],
        },
        {
            "kind": "figure",
            "title": "训练闭环",
            "figure_index": 0,
            "caption": "模型给自己出题打分再训练，看左侧循环箭头",
        },
        {"kind": "closing", "title": "谢谢", "subtitle": "欢迎讨论"},
    ],
}

# 1x1 红点 PNG
_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
    "0000000d4944415478da63f8cfc0f01f0005050102cf9e3b2e0000000049454e44ae426082"
)


def test_validate_deck_spec_rules():
    spec = DeckSpec.model_validate(DECK)
    assert validate_deck_spec(spec, figure_indices={0}) == []

    bad = DeckSpec.model_validate(
        {
            "title": "t",
            "slides": [
                {
                    "kind": "cover",
                    "title": "一个特别长长长长长长长长长长长长长长的标题——还带破折号",
                },
                {"kind": "content", "title": "要点", "bullets": ["条目" + "很长" * 30]},
                {"kind": "figure", "title": "图", "figure_index": 9},
            ],
        }
    )
    errors = validate_deck_spec(bad, figure_indices={0})
    joined = "\n".join(errors)
    assert "破折号" in joined and "上限" in joined
    assert "figure_index 非法" in joined and "讲解" in joined


def test_build_deck_enforces_template_rules():
    spec = DeckSpec.model_validate(DECK)
    data = build_deck(spec, {0: _PNG})
    assert len(data) < 4_000_000  # 模板样例页与大媒体未混入产物
    prs = PptxFile(io.BytesIO(data))
    slides = list(prs.slides)
    assert len(slides) == 5
    # 标题 30pt、正文 18pt、二级要点 16pt
    sizes: set[float] = set()
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            sizes.add(run.font.size.pt)
    assert {30.0, 18.0, 16.0} <= sizes
    # figure 页图片已插入（占位符内嵌图：元素树里出现 blip 引用）
    assert any(sh._element.xpath('.//*[local-name()="blip"]') for sh in slides[3].shapes)  # noqa: SLF001


SKILL_MD = """---
name: 论文分享 PPT 制作
description: 分享用 PPT 的模板规范
---

# 论文分享 PPT 制作

标题短一点，正文一行一个短句。
"""


async def _setup(client):
    token = await register_and_login(client)
    headers = {"Authorization": f"Bearer {token}"}
    resp = await client.post("/api/projects", json={"name": "ppt-proj"}, headers=headers)
    return headers, resp.json()["id"]


async def test_import_skill_md(client):
    headers, _ = await _setup(client)
    resp = await client.post(
        "/api/skills/import-md",
        json={"content": SKILL_MD, "targets": ["present.slides", "present.outline"]},
        headers=headers,
    )
    assert resp.status_code == 201, resp.text
    skill = resp.json()
    assert skill["name"] == "论文分享 PPT 制作"
    assert skill["kind"] == "guidance"
    assert skill["current_version"]["manifest"]["targets"] == [
        "present.slides",
        "present.outline",
    ]
    assert "标题短一点" in skill["current_version"]["body"]
    # 再导入一次：slug 自动加后缀
    resp = await client.post(
        "/api/skills/import-md",
        json={"content": SKILL_MD, "targets": ["present.slides"]},
        headers=headers,
    )
    assert resp.status_code == 201
    assert resp.json()["slug"] != skill["slug"]


async def test_create_presentation_voyage(client, queue_stub):
    headers, project_id = await _setup(client)
    # 无论文 → 404；papers 属于项目校验
    resp = await client.post(
        f"/api/projects/{project_id}/presentations",
        json={"paper_ids": [str(uuid.uuid4())], "mode": "single"},
        headers=headers,
    )
    assert resp.status_code == 404

    from app.core.db import get_sessionmaker

    async with get_sessionmaker()() as session:
        paper = await add_paper(session,
            project_id=uuid.UUID(project_id),
            title="Self-Rewarding Language Models",
            abstract="LLM as its own reward model.",
            status="compiled",
            source="manual",
        )
        session.add(paper)
        await session.commit()
        paper_id = str(paper.id)

    resp = await client.post(
        f"/api/projects/{project_id}/presentations",
        json={"paper_ids": [paper_id], "mode": "single", "notes": "面向组会"},
        headers=headers,
    )
    assert resp.status_code == 201, resp.text
    voyage = resp.json()
    assert voyage["kind"] == "presentation"
    assert ("run_voyage", (voyage["id"],), {}) in queue_stub.jobs
    # 文件未生成时下载 404
    resp = await client.get(f"/api/presentations/{voyage['id']}/file", headers=headers)
    assert resp.status_code == 404
    # survey 需要多篇由 mode 校验（single 传两篇 → 422）
    resp = await client.post(
        f"/api/projects/{project_id}/presentations",
        json={"paper_ids": [paper_id, paper_id], "mode": "single"},
        headers=headers,
    )
    assert resp.status_code == 422


# ---- 视觉审查必须覆盖整份 deck（#436）----


class _RecordingVLM:
    """记录每次视觉审查收到几张图、提示词声称的是哪几页。

    第一轮全部回 ok，用来验证覆盖范围；把 ``fail_page`` 设成某个绝对页码，
    那一批会回一条问题，用来验证页码基准。
    """

    def __init__(self, fail_page: int | None = None) -> None:
        self.calls: list[dict] = []
        self.fix_prompts: list[str] = []
        self.deck: dict = {}  # 由 _run_build 填成被测 deck，修复调用原样回它
        self.fail_page = fail_page

    async def complete(self, stage, messages, *, images=None, **kwargs):
        import json as _json
        import re as _re

        from app.core.llm.base import CompletionResult

        user = messages[-1].content
        if not images:
            # 没带图 = 修复调用。deck 本身是合规的，所以只该由视觉问题触发；
            # 没报问题却走到这里说明文本校验规则变了，此时"审了几页"的断言已不成立。
            assert self.fail_page is not None, "deck 应当已合规，不该触发文本修复"
            self.fix_prompts.append(user)
            return CompletionResult(content=_json.dumps(self.deck), model="fake-llm")

        span = _re.search(r"第 (\d+)-(\d+) 页", user)
        first, last = (int(span.group(1)), int(span.group(2))) if span else (0, 0)
        self.calls.append({"images": len(images or []), "first": first, "last": last, "user": user})

        issues = []
        if self.fail_page is not None and first <= self.fail_page <= last:
            issues = [{"slide": self.fail_page, "problem": "文字溢出", "fix": "缩短标题"}]
        return CompletionResult(
            content=_json.dumps({"ok": not issues, "issues": issues}), model="fake-vlm"
        )


def _wide_deck(pages: int) -> dict:
    """一份 pages 页的合法 deck（首页封面、末页结语，中间填内容页）。"""
    slides = [DECK["slides"][0], DECK["slides"][1]]
    slides += [
        {"kind": "content", "title": f"第 {i} 节", "bullets": [f"要点 {i}", f"补充 {i}"]}
        for i in range(len(slides) + 1, pages)
    ]
    slides.append(DECK["slides"][-1])
    return {"title": "长 deck", "slides": slides}


async def _run_build(monkeypatch, *, pages: int, vlm: _RecordingVLM):
    """跑 present.build，把渲染与 soffice 探测都替换掉（CI 没有 soffice）。"""
    from app.agents.voyage import actions_present
    from app.agents.voyage.actions import ActionContext
    from app.models.voyage import VoyageRun

    monkeypatch.setattr(actions_present, "soffice_available", lambda: True)
    monkeypatch.setattr(actions_present, "render_slide_images", lambda _pptx: [_PNG] * pages)

    async def _no_papers(_ctx):
        return []

    monkeypatch.setattr(actions_present, "_load_papers", _no_papers)

    vlm.deck = _wide_deck(pages)
    run = VoyageRun(id=uuid.uuid4(), kind="presentation", goal="测试")
    ctx = ActionContext(run=run, llm=vlm, checkpoint={"present_deck": vlm.deck})
    return await actions_present.present_build(ctx, {})


def test_render_page_cap_matches_deck_cap():
    """渲染上限必须等于 deck 页数上限，否则超出的页连图都没有，更谈不上被审。

    #436 报的是 images[:8]，但其实有两道截断：render_slide_images 的 max_pages 曾默认 12，
    而 DeckSpec 允许 25 页——13-25 页从来没被渲染过，也从来没人发现。两个上限分处两个文件、
    各写各的字面量，这条用例就是不让它们再各自漂。
    """
    import inspect

    from app.services.presentation import MAX_DECK_SLIDES, DeckSpec, render_slide_images

    default = inspect.signature(render_slide_images).parameters["max_pages"].default
    assert default == MAX_DECK_SLIDES, f"渲染上限 {default} ≠ deck 上限 {MAX_DECK_SLIDES}"

    # 且 MAX_DECK_SLIDES 确实是 deck 真正能达到的页数（常量与校验规则也得对得上）
    DeckSpec.model_validate(_wide_deck(MAX_DECK_SLIDES))
    with pytest.raises(ValidationError):
        DeckSpec.model_validate(_wide_deck(MAX_DECK_SLIDES + 1))


async def test_visual_review_covers_every_rendered_page(monkeypatch, tmp_path):
    """25 页的 deck，每一页都要恰好被送审一次。

    以前提示词说「共 N 页」却只发 images[:8]：第 9 页往后从未被看过，而模型对收到的
    8 张回 ok=true，循环就退出，对外表现为整份 deck 视觉审查通过。这种失败没有任何
    报错，只能靠读代码发现——所以这里直接钉住覆盖范围。
    """
    from app.core.config import get_settings

    monkeypatch.setattr(get_settings(), "data_dir", str(tmp_path), raising=False)
    vlm = _RecordingVLM()
    result = await _run_build(monkeypatch, pages=25, vlm=vlm)

    covered: list[int] = []
    for call in vlm.calls:
        covered.extend(range(call["first"], call["last"] + 1))
        # 提示词声称的页数必须等于真正附上的图片数——这两者不一致正是 #436 的成因
        assert call["last"] - call["first"] + 1 == call["images"], call["user"]

    assert covered == list(range(1, 26)), f"每页恰好一次，实际：{covered}"
    assert result["rendered_pages"] == 25
    assert result["reviewed_pages"] == 25


async def test_visual_issue_pages_are_absolute(monkeypatch, tmp_path):
    """第二批里报的问题，页码要是整份 deck 的绝对页码，不是批内序号。

    不要求绝对页码的话，每批都会从 1 开始编号，修复时会照着改错页。
    """
    from app.core.config import get_settings

    monkeypatch.setattr(get_settings(), "data_dir", str(tmp_path), raising=False)
    vlm = _RecordingVLM(fail_page=9)  # 第二批的第一页
    result = await _run_build(monkeypatch, pages=25, vlm=vlm)

    assert result["visual_fix_rounds"] >= 1, "报了问题就该进修复轮"

    # 第 9 页落在第二批（9-16）的第一张。诊断必须说"第 9 页"——若按批内序号编号就成了"第 1 页"，
    # 修复会照着改封面。
    assert vlm.fix_prompts, "视觉问题应当触发一次修复调用"
    assert "第 9 页" in vlm.fix_prompts[0], vlm.fix_prompts[0]

    # 且送审提示词本身要向模型交代绝对页码基准，否则上面那个页码只是碰巧对
    second = vlm.calls[1]
    assert (second["first"], second["last"]) == (9, 16), vlm.calls
    assert "绝对页码" in second["user"], second["user"]
    pages = [i["slide"] for i in result["last_visual_issues"]]
    assert pages == [9], f"应当是绝对页码 9，实际 {pages}"
